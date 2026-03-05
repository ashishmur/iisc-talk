#!/usr/bin/env python3
"""Generate slide deck: AI & The New Management Playbook — IISc DoMS, March 6, 2026"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ─── COLOR PALETTE ───────────────────────────────────────

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = {
    'dark_bg':    rgb('#0D1B2A'), 'dark_srf':   rgb('#1B2838'),
    'white':      rgb('#FFFFFF'), 'near_black': rgb('#242424'),
    'med_gray':   rgb('#616161'), 'light_gray': rgb('#F5F5F5'),
    'border':     rgb('#E0E0E0'), 'accent':     rgb('#0F6CBD'),
    'accent_hov': rgb('#115EA3'), 'accent_lt':  rgb('#E8F4FD'),
    'teal':       rgb('#00B7C3'), 'gold':       rgb('#FFB900'),
    'green':      rgb('#107C10'), 'red':        rgb('#D13438'),
    'purple':     rgb('#8764B8'), 'orange':     rgb('#F7630C'),
    'dark_fg':    rgb('#FFFFFF'), 'dark_fg2':   rgb('#B0BEC5'),
}
FONT = 'Segoe UI'

# ─── HELPERS ─────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, w, h, text, sz=14, bold=False, color=None, align='left', lsp=None):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    amap = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.name = FONT
        p.font.color.rgb = color or C['near_black']
        p.alignment = amap.get(align, PP_ALIGN.LEFT)
        if lsp:
            p.space_after = Pt(lsp)
    return txbox

def add_line(slide, left, top, w, color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or C['accent']
    shape.line.fill.background()

def add_rect(slide, left, top, w, h, fill_color, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(st, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_bullets(slide, left, top, w, h, items, sz=18, color=None, spacing=6):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.name = FONT
        p.font.color.rgb = color or C['near_black']
        p.space_before = Pt(0)
        p.space_after = Pt(spacing)

def add_card(slide, left, top, w, h, title, bullets, accent=None, bg=None):
    accent = accent or C['accent']
    add_rect(slide, left, top, w, h, bg or C['light_gray'], rounded=True)
    add_rect(slide, left, top, w, 0.06, accent)
    add_text(slide, left+0.2, top+0.2, w-0.4, 0.4, title, sz=16, bold=True, color=accent)
    if bullets:
        add_bullets(slide, left+0.2, top+0.65, w-0.4, h-0.85, bullets, sz=13, spacing=4)

def add_takeaway(slide, left, top, w, text, accent=None):
    accent = accent or C['accent']
    add_rect(slide, left, top, w, 0.65, C['accent_lt'])
    add_rect(slide, left, top, 0.06, 0.65, accent)
    add_text(slide, left+0.2, top+0.1, w-0.4, 0.45, text, sz=14, bold=True, color=C['accent_hov'])

def add_logo(slide):
    logo = Path('~/.claude/assets/microsoft-logo.png').expanduser()
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(11.5), Inches(7.0), width=Inches(0.9))


# ─── PRESENTATION SETUP ─────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ═════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_line(s, 2.0, 2.5, 9.333)
add_text(s, 1.0, 2.7, 11.333, 1.2, "AI & The New Management Playbook",
         sz=42, bold=True, color=C['dark_fg'], align='center')
add_text(s, 1.0, 3.9, 11.333, 0.5, "Ashish Muralidharan  |  Product Manager, Microsoft",
         sz=20, color=C['accent'], align='center')
add_text(s, 1.0, 4.5, 11.333, 0.5, "IISc Department of Management Studies  |  March 6, 2026",
         sz=16, color=C['dark_fg2'], align='center')
add_line(s, 2.0, 5.2, 9.333)
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 2: THE ANXIETY POLL
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.5, 2.5, 10.333, 2.5,
    "How many of you feel anxious\nabout your future\nbecause of AI?",
    sz=36, bold=True, color=C['dark_fg'], align='center', lsp=16)
add_text(s, 1.5, 5.5, 10.333, 0.5, "Show of hands.",
    sz=18, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 3: VUCA WORLD
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7,
    "We live in a VUCA world. It's not cyclical. It's a one-way street.", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

vuca = [
    ("V", "Volatility", "Rapid, unpredictable change", C['accent']),
    ("U", "Uncertainty", "Lack of predictability", C['teal']),
    ("C", "Complexity", "Interconnected forces", C['purple']),
    ("A", "Ambiguity", "Unclear cause and effect", C['orange']),
]
for i, (letter, word, desc, clr) in enumerate(vuca):
    x = 0.5 + i * 3.1
    add_rect(s, x, 1.8, 2.9, 3.0, C['light_gray'], rounded=True)
    add_rect(s, x, 1.8, 2.9, 0.06, clr)
    add_text(s, x+0.2, 2.1, 2.5, 0.8, letter, sz=48, bold=True, color=clr)
    add_text(s, x+0.2, 2.9, 2.5, 0.4, word, sz=20, bold=True)
    add_text(s, x+0.2, 3.4, 2.5, 0.8, desc, sz=14, color=C['med_gray'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "Each technology disruption resets the VUCA cycle. AI is the latest, and possibly largest, reset.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 4: TECHNOLOGY CYCLES
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7,
    "Technology has always been at the center of wealth creation", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

cycles = [
    ("Shipping", "Colonial expansion,\nglobal trade routes", "1500s-1800s", C['accent']),
    ("Industrial\nRevolution", "Mechanized labor,\nrestructured society", "1800s-1900s", C['teal']),
    ("Information\nTechnology", "Internet, tech giants,\ndigital economy", "1970s-2020s", C['purple']),
    ("Artificial\nIntelligence", "The next reset.\nWe are here.", "2020s onward", C['gold']),
]
for i, (era, desc, period, clr) in enumerate(cycles):
    x = 0.5 + i * 3.1
    add_circle(s, x+1.2, 1.8, 0.5, clr)
    add_text(s, x, 2.5, 2.9, 0.8, era, sz=18, bold=True, align='center')
    add_text(s, x, 3.4, 2.9, 0.8, desc, sz=13, color=C['med_gray'], align='center')
    add_text(s, x, 4.3, 2.9, 0.4, period, sz=12, bold=True, color=clr, align='center')
    if i < 3:
        add_text(s, x+2.7, 1.9, 0.5, 0.4, "\u2192", sz=18, color=C['border'], align='center')

add_takeaway(s, 0.5, 5.8, 12.333,
    "The largest wealth-creating firms in every era are technology companies. All the billionaires come from there.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 5: ECONOMIC EQUATION
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 0.5, 0.5, 12.333, 0.7, "The fundamental economic equation",
    sz=28, bold=True, color=C['dark_fg'])
add_line(s, 0.5, 1.3, 5.0)

factors = ["Land", "Labor", "Capital", "Technology"]
factor_colors = [C['dark_fg2'], C['dark_fg2'], C['dark_fg2'], C['accent']]
for i, (f, fc) in enumerate(zip(factors, factor_colors)):
    x = 1.0 + i * 3.0
    add_text(s, x, 2.5, 2.5, 0.8, f, sz=36, bold=True, color=fc, align='center')
    if i < 3:
        add_text(s, x+2.3, 2.5, 0.5, 0.8, "+", sz=36, color=C['dark_fg2'], align='center')

add_text(s, 1.0, 4.0, 11.333, 0.6,
    "Technology is fundamentally replacing labor in this equation.",
    sz=24, color=C['gold'], align='center')
add_text(s, 1.0, 5.0, 11.333, 0.5,
    "We are at the outset of this shift. And it changes everything.",
    sz=16, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 6: SECTION — THE PATTERN
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 2.0, 11.333, 0.8, "02", sz=48, bold=True, color=C['accent'])
add_line(s, 1.0, 2.9, 4.0)
add_text(s, 1.0, 3.1, 11.333, 0.8, "The Pattern", sz=36, bold=True, color=C['dark_fg'])
add_text(s, 1.0, 3.9, 11.333, 0.5, "Technology vs. labor: a history",
    sz=18, color=C['dark_fg2'])
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 7: THE 4-STEP SCRIPT
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7,
    "Every technology cycle follows the same script", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

steps = [
    ("1", "Technology arrives", "Anxiety about jobs"),
    ("2", "Physical labor displaced", "People move 'up' to operational roles"),
    ("3", "Operational roles displaced", "People move 'up' to intellectual roles"),
    ("4", "New job categories appear", "Roles that never existed before"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.6 + i * 1.0
    add_circle(s, 1.0, y, 0.5, C['accent'])
    add_text(s, 1.0, y+0.05, 0.5, 0.4, num, sz=18, bold=True, color=C['white'], align='center')
    add_text(s, 1.8, y+0.05, 4.5, 0.4, title, sz=20, bold=True)
    add_text(s, 7.0, y+0.08, 5.0, 0.4, desc, sz=16, color=C['med_gray'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "The software programmer never existed until it did. Each cycle creates roles we cannot yet imagine.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 8: RUNNING OUT OF RUNGS
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "But we may be running out of rungs", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

rungs = [
    ("Physical labor", "Removed by mechanization", 10.0, C['accent'], "\u2713"),
    ("Operational labor", "Removed by information technology", 8.5, C['teal'], "\u2713"),
    ("Intellectual labor", "Being removed by AI", 7.0, C['gold'], "\u2190 We are here"),
    ("???", "What's left?", 5.0, C['red'], ""),
]
for i, (label, desc, bw, clr, status) in enumerate(rungs):
    y = 1.6 + i * 1.05
    add_rect(s, 1.0, y, bw, 0.8, clr, rounded=True)
    add_text(s, 1.2, y+0.1, bw-0.4, 0.3, label, sz=18, bold=True, color=C['white'])
    add_text(s, 1.2, y+0.4, bw-0.4, 0.3, desc, sz=12, color=C['white'])
    if status:
        add_text(s, 1.0+bw+0.2, y+0.2, 2.5, 0.4, status, sz=14, bold=True, color=clr)

add_takeaway(s, 0.5, 5.8, 12.333,
    "Beyond intellect, what can we contribute? We don't have an answer for that yet.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 9: QUOTE — CEILING
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 1.5, 1.0, 1.0, "\u201C", sz=72, bold=True, color=C['accent'])
add_text(s, 1.5, 2.3, 10.0, 2.5,
    "If you don't need intellectual labor anymore, what's left?\nYou sell your time, or you sell your capital.\nLabor has nowhere else to move.\nWe've hit the ceiling.",
    sz=26, color=C['dark_fg'], lsp=12)
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 10: ROBOTICS
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "The counter-narrative: embodied AI", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (name, desc, clr) in enumerate([
    ("Tesla Optimus", ["Humanoid general-purpose robot"], C['accent']),
    ("Unitree", ["Agile, affordable robotics at scale"], C['teal']),
    ("Hyundai Robotics", ["Industrial and service robotics"], C['purple']),
]):
    add_card(s, 0.5+i*4.11, 1.8, 3.94, 2.0, name, desc, accent=clr)

add_text(s, 0.5, 4.2, 12.333, 0.8,
    "When AI gets a body to think, see, and operate in the physical world,\nyou're talking about complete replacement of human labor. Everything.",
    sz=18, align='center')
add_takeaway(s, 0.5, 5.8, 12.333,
    "The 'golden era of hardware jobs' may also be the last frontier of human labor.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 11: 5 LEVELS OF AGI
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "OpenAI's five levels of AGI: where are we?", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

levels = [
    ("1", "Conversational AI", "ChatGPT", "Done", C['green']),
    ("2", "Reasoning", "Chain-of-thought, problem solving", "Done", C['green']),
    ("3", "Agents", "Plan, reason, act, learn", "We are here", C['gold']),
    ("4", "AI Organizations", "Multiple agents coordinating", "Emerging", C['orange']),
    ("5", "AGI", "General intelligence, all domains", "Unknown", C['red']),
]
# Header
cols_x = [0.5, 1.5, 4.0, 8.5]
cols_w = [1.0, 2.5, 4.5, 3.5]
for j, (header, x, w) in enumerate(zip(["Level", "Capability", "Description", "Status"], cols_x, cols_w)):
    add_rect(s, x, 1.5, w, 0.5, C['accent'])
    add_text(s, x+0.1, 1.55, w-0.2, 0.4, header, sz=14, bold=True, color=C['white'])

for i, (level, name, desc, status, sclr) in enumerate(levels):
    y = 2.02 + i * 0.7
    bg = rgb('#FFF8E1') if i == 2 else (C['white'] if i % 2 == 0 else C['light_gray'])
    for j, (x, w) in enumerate(zip(cols_x, cols_w)):
        add_rect(s, x, y, w, 0.65, bg)
    vals = [level, name, desc, status]
    for j, (val, x, w) in enumerate(zip(vals, cols_x, cols_w)):
        clr = sclr if j == 3 else C['near_black']
        bld = (j == 3) or (i == 2)
        add_text(s, x+0.1, y+0.12, w-0.2, 0.4, val, sz=14, bold=bld, color=clr)

add_takeaway(s, 0.5, 5.8, 12.333,
    "We are at Level 3: Agents. Systems that plan, reason, act, and learn. This has already changed everything.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 12: AGENTS — BRANCHING FROM LEVEL 3
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Level 3: Agents", sz=28, bold=True, color=C['dark_fg'])
add_text(s, 0.5, 0.85, 4.0, 0.3, "We are here, 2025", sz=16, bold=True, color=C['gold'])
add_line(s, 0.5, 1.2, 5.0)

# Left side: Agent definition with 3 components
add_text(s, 0.5, 1.5, 5.5, 0.5,
    "What is an agent?", sz=20, bold=True, color=C['accent'])
add_text(s, 0.5, 2.0, 5.5, 0.6,
    "A system that repeatedly pursues a goal by:",
    sz=16, color=C['dark_fg'])

agent_components = [
    ("i)", "Planning and reasoning", "Break down complex goals into steps", C['accent']),
    ("ii)", "Acting on the environment", "Using tools to execute actions", C['teal']),
    ("iii)", "Context and memory", "Retaining knowledge across interactions", C['purple']),
]
for i, (num, title, desc, clr) in enumerate(agent_components):
    y = 2.7 + i * 1.0
    add_rect(s, 0.7, y, 5.1, 0.85, C['dark_srf'], rounded=True)
    add_rect(s, 0.7, y, 0.06, 0.85, clr)
    add_text(s, 0.95, y+0.08, 4.7, 0.35, f"{num}  {title}", sz=16, bold=True, color=clr)
    add_text(s, 0.95, y+0.45, 4.7, 0.35, desc, sz=13, color=C['dark_fg2'])

# Right side: Examples of AI agents
add_text(s, 6.8, 1.5, 5.5, 0.5,
    "Examples of AI agents", sz=20, bold=True, color=C['gold'])

agent_examples = [
    ("Coding Agent", "Writes, tests, and deploys code", C['accent']),
    ("Research Agent", "Finds, synthesizes, and summarizes information", C['teal']),
    ("Customer Service Agent", "Handles queries, resolves issues 24/7", C['purple']),
    ("Data Analysis Agent", "Processes data, generates insights", C['gold']),
    ("Creative Agent", "Generates content, designs, and copy", C['orange']),
]
for i, (name, desc, clr) in enumerate(agent_examples):
    y = 2.1 + i * 0.85
    add_rect(s, 7.0, y, 5.3, 0.7, C['dark_srf'], rounded=True)
    add_rect(s, 7.0, y, 0.06, 0.7, clr)
    add_text(s, 7.25, y+0.05, 4.8, 0.3, name, sz=15, bold=True, color=clr)
    add_text(s, 7.25, y+0.35, 4.8, 0.3, desc, sz=12, color=C['dark_fg2'])

add_text(s, 0.5, 5.8, 12.333, 0.5,
    "For the first time, the entire signal-to-closure loop has been completely automated.",
    sz=16, bold=True, color=C['gold'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 13: SECTION — GROUND ZERO
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 2.0, 11.333, 0.8, "03", sz=48, bold=True, color=C['accent'])
add_line(s, 1.0, 2.9, 4.0)
add_text(s, 1.0, 3.1, 11.333, 0.8, "Ground Zero: Software", sz=36, bold=True, color=C['dark_fg'])
add_text(s, 1.0, 3.9, 11.333, 0.5, "The first industry to fall", sz=18, color=C['dark_fg2'])
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 14: WHY SOFTWARE FIRST
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Why software is the first industry to fall", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (title, bullets, clr) in enumerate([
    ("100% Digital", ["All work exists as data", "Code, docs, git history"], C['accent']),
    ("Fully Closed Loop", ["Signal-to-closure is digital", "No physical component"], C['teal']),
    ("Rich Training Data", ["Decades of open-source code", "Documentation everywhere"], C['purple']),
]):
    add_card(s, 0.5+i*4.11, 1.6, 3.94, 2.5, title, bullets, accent=clr)

add_takeaway(s, 0.5, 5.8, 12.333,
    "Single engineers now spin up 20 agents that read a spec and build the entire product end-to-end.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 15: VIBE CODING
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "The new reality of building software", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

# Karpathy card
add_rect(s, 0.5, 1.6, 5.9, 2.5, C['light_gray'], rounded=True)
add_text(s, 0.7, 1.8, 5.5, 0.3, "Andrej Karpathy", sz=14, bold=True, color=C['accent'])
add_text(s, 0.7, 2.1, 5.5, 0.3, "Ex-Tesla AI, Ex-OpenAI", sz=11, color=C['med_gray'])
add_text(s, 0.7, 2.5, 5.5, 1.0,
    "Coined 'vibe coding': top developers haven't written code by hand in 100+ days.", sz=18)

# Naval card
add_rect(s, 6.6, 1.6, 5.9, 2.5, C['light_gray'], rounded=True)
add_text(s, 6.8, 1.8, 5.5, 0.3, "Naval Ravikant", sz=14, bold=True, color=C['teal'])
add_text(s, 6.8, 2.1, 5.5, 0.3, "Entrepreneur, AngelList founder", sz=11, color=C['med_gray'])
add_text(s, 6.8, 2.5, 5.5, 1.0,
    "\"Product management is vibe coding.\nEngineering is prompt engineering.\"", sz=18)

add_text(s, 0.5, 4.5, 12.333, 0.6,
    "Even playing field: no domain expertise barrier, no expensive labor needed.\nAnyone can build. The cost to produce software is approaching zero.",
    sz=18, align='center')
add_takeaway(s, 0.5, 5.8, 12.333,
    "You don't have to be literate in engineering to build a product anymore.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 16: QUOTE — COST TO ZERO
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 1.5, 1.0, 1.0, "\u201C", sz=72, bold=True, color=C['accent'])
add_text(s, 1.5, 2.5, 10.0, 2.0,
    "The cost to build software has come down to zero.\nThat's not hyperbole.\nThat's what's happening right now.",
    sz=28, color=C['dark_fg'], lsp=12)
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 17: DISTRIBUTION WINS
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7,
    "When production costs hit zero, distribution wins", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

add_card(s, 0.5, 1.6, 5.9, 3.5, "What changes", [
    "Anyone can build software",
    "No domain expertise barrier",
    "Every market becomes overcrowded",
], accent=C['red'])

add_card(s, 6.6, 1.6, 5.9, 3.5, "What matters now", [
    "Distribution and marketing",
    "Who can reach customers first?",
    "Every market becomes a red ocean",
], accent=C['green'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "\"The new job for everyone has become marketing.\" Think Porter's Five Forces when barriers to entry collapse.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 18: REAL CONSEQUENCES
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "This is already happening", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (stat, desc, clr) in enumerate([
    ("50%", "Block announced layoffs\nspecifically due to AI", C['red']),
    ("25-30%", "IT company stocks down,\nmarkets pricing in change", C['orange']),
    ("10x", "AI mentions in financial\ndisclosures have exploded", C['accent']),
]):
    x = 0.5 + i * 4.11
    add_rect(s, x, 1.6, 3.94, 3.0, C['light_gray'], rounded=True)
    add_text(s, x, 1.9, 3.94, 0.8, stat, sz=48, bold=True, color=clr, align='center')
    add_text(s, x+0.3, 2.9, 3.34, 0.8, desc, sz=15, align='center')

add_takeaway(s, 0.5, 5.8, 12.333,
    "Companies without an AI strategy are being punished by markets. Every tech company faces the innovator's dilemma.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 19: DEFLATION VS ABUNDANCE
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Two sides of the economic argument", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

add_card(s, 0.5, 1.6, 5.9, 3.5, "The abundance thesis", [
    "Remove humans from process loops",
    "Fundamental deflation in costs",
    "Goods become cheaper for everyone",
    "Era of abundance, not scarcity",
], accent=C['green'])

add_card(s, 6.6, 1.6, 5.9, 3.5, "The counter-thesis", [
    "What new jobs can appear?",
    "Beyond intellectual work, what's left?",
    "Only capital owners have a role?",
    "The honest answer: we don't know",
], accent=C['red'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "Our economics never taught us to prepare for this. You've broken the equation with a near-zero factor.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 20: AI VALUE CHAIN
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 0.5, 0.5, 12.333, 0.7, "Where value accrues in the AI stack",
    sz=28, bold=True, color=C['dark_fg'])
add_line(s, 0.5, 1.3, 5.0)

chain = ["Energy", "Chips", "Fabricators", "Frontier\nLabs", "Hyper-\nscalers", "Apps", "Distribution"]
highlights = {4, 6}  # Hyperscalers and Distribution
for i, step in enumerate(chain):
    x = 0.3 + i * 1.85
    clr = C['accent'] if i in highlights else C['dark_srf']
    add_rect(s, x, 2.2, 1.6, 1.0, clr, rounded=True)
    add_text(s, x, 2.3, 1.6, 0.8, step, sz=13, bold=True, color=C['dark_fg'], align='center')
    if i < 6:
        add_text(s, x+1.6, 2.4, 0.25, 0.5, "\u2192", sz=16, color=C['dark_fg2'], align='center')

add_text(s, 1.0, 3.8, 11.0, 0.5,
    "The model is a commodity. Value lives in compute, context, and distribution.",
    sz=22, bold=True, color=C['gold'], align='center')

for i, (title, desc) in enumerate([
    ("Compute", "Infrastructure layer"),
    ("Context", "Organizational knowledge"),
    ("Distribution", "Reaching customers"),
]):
    x = 1.5 + i * 3.8
    add_text(s, x, 4.6, 3.0, 0.4, title, sz=18, bold=True, color=C['accent'], align='center')
    add_text(s, x, 5.0, 3.0, 0.4, desc, sz=13, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 21: AGENTS DON'T NEED UI
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 0.5, 0.5, 12.333, 0.7, "Agents don't need UI",
    sz=28, bold=True, color=C['dark_fg'])
add_line(s, 0.5, 1.3, 5.0)
add_text(s, 1.0, 2.0, 11.0, 1.5,
    "UI was built for humans.\nAgents communicate through protocols.\nYour AI talks to other AIs and gets things done.",
    sz=24, color=C['dark_fg'], align='center', lsp=10)
add_text(s, 1.0, 4.0, 11.0, 0.5, "Anthropic's Model Context Protocol (MCP)",
    sz=18, bold=True, color=C['accent'], align='center')
add_text(s, 1.0, 4.6, 11.0, 0.6,
    "The application layer is under threat.\nWhy would you subject yourself to UI friction?",
    sz=18, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 22: SECTION — FROM THE INSIDE
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 2.0, 11.333, 0.8, "04", sz=48, bold=True, color=C['accent'])
add_line(s, 1.0, 2.9, 4.0)
add_text(s, 1.0, 3.1, 11.333, 0.8, "From the Inside", sz=36, bold=True, color=C['dark_fg'])
add_text(s, 1.0, 3.9, 11.333, 0.5, "What we learned at Microsoft",
    sz=18, color=C['dark_fg2'])
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 23: MICROSOFT'S AI JOURNEY
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Microsoft's AI deployment: what we learned", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

add_card(s, 0.5, 1.6, 5.9, 2.2, "What failed: top-down adoption", [
    "\"Here's how to use AI to be productive\"",
    "Too much nuance in individual work",
    "Generic training doesn't stick",
], accent=C['red'])

add_card(s, 6.6, 1.6, 5.9, 2.2, "What worked: bottom-up adoption", [
    "Unlimited compute for product managers",
    "Dedicated time to explore capabilities",
    "Incentive structures tied to AI usage",
], accent=C['green'])

add_text(s, 0.5, 4.3, 12.333, 0.6,
    "Software engineers who adopted this are seeing 20x productivity.\nThat's at Microsoft, with our own tools.",
    sz=20, bold=True, color=C['accent'], align='center')
add_takeaway(s, 0.5, 5.8, 12.333,
    "The more you use it and understand the patterns, the better you figure out what to build for customers.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 24: PROJECT NEBULA
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Project Nebula: building the muscle at scale", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (num, label, clr) in enumerate([
    ("100", "Product Managers", C['accent']),
    ("600", "Engineers", C['teal']),
    ("40-50", "Designers", C['purple']),
]):
    x = 0.5 + i * 4.11
    add_text(s, x, 1.5, 3.94, 0.8, num, sz=48, bold=True, color=clr, align='center')
    add_text(s, x, 2.2, 3.94, 0.3, label, sz=16, color=C['med_gray'], align='center')

add_bullets(s, 0.7, 2.8, 11.5, 2.5, [
    "Generate better ideas with AI",
    "Create mockups and visual prototypes",
    "Build working prototypes with models inside them",
    "Deploy AI workflows into everyday work",
    "Turn discoveries into learning modules for the org",
], sz=16, spacing=6)

add_takeaway(s, 0.5, 5.8, 12.333,
    "New joiners inherit tribal knowledge we've accumulated and ramp up from there.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 25: SECTION — THE NEW PLAYBOOK
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 2.0, 11.333, 0.8, "05", sz=48, bold=True, color=C['accent'])
add_line(s, 1.0, 2.9, 4.0)
add_text(s, 1.0, 3.1, 11.333, 0.8, "The New Playbook", sz=36, bold=True, color=C['dark_fg'])
add_text(s, 1.0, 3.9, 11.333, 0.5, "Designing for AI organizations",
    sz=18, color=C['dark_fg2'])
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 26: AGENTS VS HUMANS
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Agents have different constraints than humans", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

add_card(s, 0.5, 1.6, 5.9, 3.5, "Weaknesses", [
    "Hallucination and inconsistency",
    "Jagged intelligence across domains",
    "Limited domain expertise",
    "Memory and context challenges",
], accent=C['red'])

add_card(s, 6.6, 1.6, 5.9, 3.5, "Strengths", [
    "Tireless and persistent",
    "Can reference massive scope",
    "Plan and reason with available tools",
    "Operate 24/7 without burnout",
], accent=C['green'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "Managers of the future will design workflows for AI agents the way managers of the past designed them for humans.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 27: PM WORKFLOW → AGENT WORKFLOW
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7,
    "Organizational design principles apply to AI systems", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

pm_steps = ["Strategy", "User\nResearch", "PRD", "Design", "Tech\nDesign", "Jira\nTickets", "Engineers\nExecute"]
for i, step in enumerate(pm_steps):
    x = 0.3 + i * 1.85
    add_rect(s, x, 1.5, 1.6, 0.8, C['accent'], rounded=True)
    add_text(s, x, 1.55, 1.6, 0.7, step, sz=11, bold=True, color=C['white'], align='center')
    if i < 6:
        add_text(s, x+1.6, 1.6, 0.25, 0.5, "\u2192", sz=14, color=C['med_gray'], align='center')

add_text(s, 0.5, 2.5, 12.333, 0.3, "Traditional PM workflow (humans at every step)",
    sz=12, color=C['med_gray'])

agent_steps = ["Strategy", "User\nResearch", "PRD", "Design", "Tech\nDesign", "Agent\nTasks", "Agents\nExecute"]
for i, step in enumerate(agent_steps):
    x = 0.3 + i * 1.85
    add_rect(s, x, 3.2, 1.6, 0.8, C['teal'], rounded=True)
    add_text(s, x, 3.25, 1.6, 0.7, step, sz=11, bold=True, color=C['white'], align='center')
    if i < 6:
        add_text(s, x+1.6, 3.3, 0.25, 0.5, "\u2192", sz=14, color=C['med_gray'], align='center')

add_text(s, 0.5, 4.2, 12.333, 0.3, "Agent workflow (same structure, AI at every step)",
    sz=12, color=C['med_gray'])
add_text(s, 0.5, 4.7, 12.333, 0.4,
    "Git for coordination. Unit tests for quality. The patterns are identical.",
    sz=18, align='center')
add_takeaway(s, 0.5, 5.8, 12.333,
    "Your ability to replace labor with compute and design value-creating workflows is the new core skill.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 28: TWO CORE SKILLS
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Two skills every manager needs now", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

add_card(s, 0.5, 1.6, 5.9, 3.5, "Prompting", [
    "The ability to instruct AI precisely",
    "Clear, structured instructions",
    "Context design and curation",
    "From vague intent to specific output",
], accent=C['accent'])

add_card(s, 6.6, 1.6, 5.9, 3.5, "Evaluations", [
    "Measuring if AI is doing a good job",
    "Think: unit tests for intelligence",
    "Defining quality at every step",
    "Feedback loops that improve output",
], accent=C['teal'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "Prompting is the new management. Evaluations are the new quality assurance.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 29: AGE OF THE POLYMATH
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 0.5, 0.5, 12.333, 0.7, "The age of the polymath",
    sz=28, bold=True, color=C['dark_fg'])
add_line(s, 0.5, 1.3, 5.0)

for i, (skill, desc, clr) in enumerate([
    ("Design", "Polished visual prototypes", C['accent']),
    ("Code", "Working systems and apps", C['teal']),
    ("Analyze", "Finance, research, strategy", C['purple']),
    ("Ship", "Deploy, test, iterate", C['gold']),
]):
    x = 0.5 + i * 3.1
    add_rect(s, x, 2.0, 2.9, 1.8, C['dark_srf'], rounded=True)
    add_rect(s, x, 2.0, 2.9, 0.06, clr)
    add_text(s, x+0.2, 2.3, 2.5, 0.4, skill, sz=22, bold=True, color=clr)
    add_text(s, x+0.2, 2.8, 2.5, 0.6, desc, sz=14, color=C['dark_fg2'])

add_text(s, 1.0, 4.3, 11.0, 0.6,
    "\"If you cannot prototype it, that idea cannot be built.\"\nThe one-person billion-dollar startup isn't a meme anymore.",
    sz=20, color=C['gold'], align='center')
add_text(s, 1.0, 5.3, 11.0, 0.6,
    "A Renaissance of thinkers who coordinate across finance, marketing,\nengineering, and growth, all at the same time.",
    sz=16, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 30: WHAT TO DO NOW
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "What you should do right now", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (num, title, desc, clr) in enumerate([
    ("1", "Build a personal AI budget", "Invest monthly. Rotate, explore, keep what works.", C['accent']),
    ("2", "Build things in your domain", "Show AI fluency through working prototypes.", C['teal']),
    ("3", "Think in workflows, not tools", "Design end-to-end processes. Build context.", C['purple']),
    ("4", "Become a generalist", "Prototype, design, code, analyze, ship.", C['gold']),
]):
    y = 1.5 + i * 1.05
    add_circle(s, 0.8, y+0.05, 0.5, clr)
    add_text(s, 0.8, y+0.1, 0.5, 0.4, num, sz=18, bold=True, color=C['white'], align='center')
    add_text(s, 1.6, y+0.05, 4.5, 0.4, title, sz=20, bold=True)
    add_text(s, 6.5, y+0.08, 6.0, 0.4, desc, sz=15, color=C['med_gray'])

add_takeaway(s, 0.5, 5.8, 12.333,
    "The most valuable managers will be the ones who figure out how to diffuse AI throughout the company.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 31: BUILD IN YOUR DOMAIN
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['white'])
add_text(s, 0.5, 0.3, 12.333, 0.7, "Show, don't tell: build in your domain", sz=28, bold=True)
add_line(s, 0.5, 1.1, 12.333)

for i, (domain, items, clr) in enumerate([
    ("Finance", ["Trading bots with technical analysis", "Automated opportunity identification", "Simulation and scoring systems"], C['accent']),
    ("Marketing", ["Brand concept generators", "Social media agents for content", "Memetic opportunity detection"], C['teal']),
    ("Growth", ["Fake-door websites to test demand", "Market validation experiments", "End-to-end campaign automation"], C['purple']),
]):
    add_card(s, 0.5+i*4.11, 1.6, 3.94, 3.0, domain, items, accent=clr)

add_takeaway(s, 0.5, 5.8, 12.333,
    "Walk into your interview showing a working AI-powered system. That's why they'll hire you, not your degree.")
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 32: RETURN TO ANXIETY
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 1.2, 11.0, 0.5, "Let's come back to where we started.",
    sz=22, color=C['dark_fg2'], align='center')
add_text(s, 1.0, 2.0, 11.0, 2.5,
    "That anxiety is the same anxiety people felt when the steam engine arrived.\nWhen the assembly line was invented. When the internet showed up.\n\nEvery single time, the people who leaned into the change,\nwho learned the new tools, who built instead of froze...\nthey didn't just survive. They defined the next era.",
    sz=20, color=C['dark_fg'], align='center', lsp=8)
add_text(s, 1.0, 5.0, 11.0, 0.8,
    "You are at one of the finest academic institutions in this country.\nYou can build a product tonight. Prototype a business idea by tomorrow morning.",
    sz=16, color=C['gold'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 33: THE REFRAME
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_text(s, 1.0, 2.0, 11.0, 2.0,
    "The ceiling on intellectual labor\nfeels like a wall if you're standing under it.\n\nBut it's a launchpad\nif you're the one building the systems\nthat sit on top of it.",
    sz=28, bold=True, color=C['dark_fg'], align='center', lsp=10)
add_text(s, 1.0, 4.8, 11.0, 0.5,
    "So let me ask again: how many of you feel anxious?",
    sz=22, color=C['accent'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SLIDE 34: CLOSING
# ═════════════════════════════════════════════════════════
s = new_slide(prs); set_bg(s, C['dark_bg'])
add_line(s, 3.0, 2.8, 7.333)
add_text(s, 1.0, 3.0, 11.333, 1.0, "Go build something.",
    sz=40, bold=True, color=C['dark_fg'], align='center')
add_line(s, 3.0, 4.2, 7.333)
add_text(s, 1.0, 5.0, 11.333, 0.5,
    "Ashish Muralidharan  |  ashish.muralidharan@microsoft.com",
    sz=16, color=C['dark_fg2'], align='center')
add_logo(s)


# ═════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════
output = Path("~/claude_personal/iisc-talk/iisc-talk-deck.pptx").expanduser()
prs.save(str(output))
print(f"Deck saved to: {output}")
print(f"Total slides: {len(prs.slides)}")
